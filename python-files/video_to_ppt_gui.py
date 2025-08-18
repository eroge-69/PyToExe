# video_to_ppt_gui.py
import os
import sys
from tkinter import Tk, filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches
from moviepy.editor import VideoFileClip
from PIL import Image

MAX_SIZE_MB = 300

def select_video():
    filetypes = [('MP4 Video', '*.mp4')]
    video_path = filedialog.askopenfilename(title="Video Seç", filetypes=filetypes)
    if not video_path:
        return None
    size_mb = os.path.getsize(video_path) / (1024*1024)
    if size_mb > MAX_SIZE_MB:
        messagebox.showerror("Hata", f"Video boyutu {MAX_SIZE_MB} MB sınırını aşıyor.")
        return None
    return video_path

def video_to_pptx(video_path):
    try:
        clip = VideoFileClip(video_path)
        duration = clip.duration
        fps = 1  # 1 kare / saniye
        frames = [clip.get_frame(t) for t in range(int(duration))]
        # Sunum oluştur
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        for idx, frame in enumerate(frames):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # boş slayt
            img = Image.fromarray(frame)
            tmp_path = os.path.join(os.environ.get('TEMP', '.'), f"_frame_{idx}.jpg")
            img.save(tmp_path, "JPEG")
            slide.shapes.add_picture(tmp_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.remove(tmp_path)
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        output_path = os.path.join(desktop, "output.pptx")
        prs.save(output_path)
        messagebox.showinfo("Başarılı", f"PPTX başarıyla oluşturuldu!\n{output_path}")
    except Exception as e:
        messagebox.showerror("Hata", f"PPTX oluşturulurken bir hata oluştu:\n{e}")

def main():
    root = Tk()
    root.withdraw()  # gizle ana pencere
    video = select_video()
    if video:
        video_to_pptx(video)
    root.destroy()

if __name__ == "__main__":
    main()
