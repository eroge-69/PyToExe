import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
import os

# =========================
#  To'liq mapping (Lotin -> Kirill va Kirill -> Lotin)
#  mappinglar kichik-cho'kin holatda saqlanadi; case saqlash logikasi pastda bor.
# =========================

# barcha apostrof/tutuq belgilar variantlari: ' ` ’ ‘ ʼ ʹ
APOS_VARIANTS = ["'", "`", "’", "‘", "ʼ", "ʹ"]

# Lotin -> Kirill (har bir kombinatsiyani pastki registrda qo'ydik)
latin_to_cyr = {
    # 2-3 harf kombinatsiyalar birinchi (uzunroq birinchi tekshiriladi)
    "sh": "ш", "ch": "ч", "ng": "нг",
    "ya": "я", "yo": "ё", "yu": "ю", "ye": "е",
    # apostrofli kombinatsiyalar (turli apostroflarni qamrab olamiz)
    "o'": "ў", "o`": "ў", "o’": "ў", "o‘": "ў", "oʼ": "ў",
    "g'": "ғ", "g`": "ғ", "g’": "ғ", "g‘": "ғ", "gʼ": "ғ",
    # alohida harflar
    "a": "а", "b": "б", "d": "д", "e": "е", "f": "ф", "g": "г",
    "h": "ҳ", "i": "и", "j": "ж", "k": "к", "l": "л", "m": "м",
    "n": "н", "o": "о", "p": "п", "q": "қ", "r": "р", "s": "с",
    "t": "т", "u": "у", "v": "в", "x": "х", "y": "й", "z": "з",
    # ayrim lotin variantlar (qo'shimcha)
    "ʼ": "ъ",  # modifier apostrophe -> hard sign
}

# Kirill -> Lotin
cyr_to_latin = {
    # uzunroq birinchi
    "ш": "sh", "ч": "ch", "нг": "ng",
    "я": "ya", "ё": "yo", "ю": "yu", "е": "e",
    "ў": "o'", "ғ": "g'",
    "а": "a", "б": "b", "д": "d", "е": "e", "ф": "f", "г": "g",
    "ҳ": "h", "и": "i", "ж": "j", "к": "k", "л": "l", "м": "m",
    "н": "n", "о": "o", "п": "p", "қ": "q", "р": "r", "с": "s",
    "т": "t", "у": "u", "в": "v", "х": "x", "й": "y", "з": "z",
    "ь": "", "ъ": "'"
}

# Maksimal kalit uzunligini hisoblaymiz
MAX_LATIN_KEY = max(len(k) for k in latin_to_cyr.keys())
MAX_CYR_KEY = max(len(k) for k in cyr_to_latin.keys())

# =========================
#  Helper: case-saqlash transliteration
# =========================

def apply_case(src_fragment: str, mapped: str) -> str:
    """
    src_fragment - asl matn fragmenti (o'z holida: katta/kichik aralash bo'lishi mumkin)
    mapped - past registersdagi xaritadan olgan natija (odatda pastki registr)
    natija: asl fragment holatini saqlashga urinar ekan
    """
    # Agar butun fragment bosh harflar bilan (masalan "SH", "O'") bo'lsa -> natijani ham upper qiling
    if src_fragment.isupper():
        return mapped.upper()
    # Agar faqat birinchi harf katta bo'lsa -> mapped ni ham birinchi harf katta qiladi
    if src_fragment[0].isupper():
        # mapped ga qarab birinchi harfni katta qilamiz; agar mapped bir nechta belgidan iborat bo'lsa,
        # faqat birinchi belgini katta qilamiz.
        return mapped[0].upper() + mapped[1:]
    # boshqacha holatda past registrni qaytaramiz
    return mapped

# =========================
#  Transliteration funksiyalari
# =========================

def transliterate_latin_to_cyrillic(text: str) -> str:
    """
    Lotin matnni Kirillga o'giradi. Turli apostroflarni qo'llab-quvvatlaydi.
    Case saqlanadi.
    """
    result = []
    i = 0
    L = len(text)
    while i < L:
        matched = False
        # sinov uchun uzunroq klavishlardan boshlaymiz
        for length in range(MAX_LATIN_KEY, 0, -1):
            if i + length > L:
                continue
            fragment = text[i:i+length]
            lower_frag = fragment.lower()
            # mapda qidiring — lekin apostrof variantlari uchun ham tekshiring
            if lower_frag in latin_to_cyr:
                mapped = latin_to_cyr[lower_frag]
                mapped_with_case = apply_case(fragment, mapped)
                result.append(mapped_with_case)
                i += length
                matched = True
                break
            # agar fragment ichida apostrof belgisi bo'lsa va pastdagi xaritada apostrof bilan berilgan ko'rinmasa,
            # har xil apostroflarni normalizatsiya qilish o'rnida mappingda barcha variantlar mavjudligini tekshiramiz
            # (keyinchalik mapga barcha apostrof variantlari qo'shilgan)
        if not matched:
            # agar mos kelmas ekan oddiy belgini qo'shamiz
            result.append(text[i])
            i += 1
    return "".join(result)

def transliterate_cyrillic_to_latin(text: str) -> str:
    """
    Kirill matnni Lotinga o'giradi. Case saqlaydi.
    """
    result = []
    i = 0
    L = len(text)
    while i < L:
        matched = False
        for length in range(MAX_CYR_KEY, 0, -1):
            if i + length > L:
                continue
            fragment = text[i:i+length]
            # NOTE: cyr_to_latin keys are in lower-case; but cyrillic letters don't have ASCII case same as Latin.
            # We check fragment as-is (since Cyrillic case also supported by .islower/.isupper)
            lower_frag = fragment.lower()
            if lower_frag in cyr_to_latin:
                mapped = cyr_to_latin[lower_frag]
                mapped_with_case = apply_case(fragment, mapped)
                result.append(mapped_with_case)
                i += length
                matched = True
                break
        if not matched:
            result.append(text[i])
            i += 1
    return "".join(result)

# =========================
#  PPTX qayta ishlash (slayd, jadval, shape, XML <a:t>)
# =========================

def process_pptx(input_path: str, output_path: str, direction: str):
    """
    direction: "latin_to_cyr" yoki "cyr_to_latin"
    """
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            # oddiy text frame va runlar
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original = run.text or ""
                        if direction == "latin_to_cyr":
                            run.text = transliterate_latin_to_cyrillic(original)
                        else:
                            run.text = transliterate_cyrillic_to_latin(original)

            # jadval ichidagi matn
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        original = cell.text or ""
                        if direction == "latin_to_cyr":
                            cell.text = transliterate_latin_to_cyrillic(original)
                        else:
                            cell.text = transliterate_cyrillic_to_latin(original)

            # XML ichidagi <a:t> elementlari (WordArt/SmartArt/Diagramma va boshqa holatlar)
            try:
                for t in shape.element.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
                    if t.text:
                        if direction == "latin_to_cyr":
                            t.text = transliterate_latin_to_cyrillic(t.text)
                        else:
                            t.text = transliterate_cyrillic_to_latin(t.text)
            except Exception:
                # agar shape.element yo'q yoki boshqa muammo bo'lsa, o'tib ketamiz
                pass

    prs.save(output_path)

# =========================
#  GUI (tkinter) — dizayn bilan, oynani rozvernutni (maximize) o'chirish
# =========================

root = tk.Tk()
root.title("PowerPoint Lotin ↔ Кирилл Tarjimon")
root.geometry("720x400")
root.resizable(False, False)
root.configure(bg="#1e1e1e")  # dark background

# Card frame with dark style
shadow = tk.Frame(root, bg="#000000")
shadow.place(relx=0.5, rely=0.5, anchor="center", width=680, height=320)
card = tk.Frame(root, bg="#2c2c2c")
card.place(relx=0.5, rely=0.5, anchor="center", width=660, height=300)

# Title
tk.Label(card, text="PowerPoint Lotin ↔ Кирилл Tarjimon",
         font=("Helvetica", 16, "bold"), fg="#ffffff", bg="#2c2c2c").pack(pady=(15, 10))

# Fayl tanlash
file_row = tk.Frame(card, bg="#2c2c2c")
file_row.pack(fill="x", pady=8, padx=20)
tk.Label(file_row, text="Fayl:", bg="#2c2c2c", fg="#ffffff", font=("Arial", 11)).pack(side="left")
file_var = tk.StringVar()
entry = tk.Entry(file_row, textvariable=file_var, font=("Arial", 11), bd=1, relief="solid", bg="#3a3a3a", fg="#ffffff", insertbackground="white")
entry.pack(side="left", padx=10, fill="x", expand=True)
def choose_file():
    path = filedialog.askopenfilename(title="PowerPoint fayl tanlang", filetypes=[("PowerPoint files", "*.pptx;*.ppt")])
    if path: file_var.set(path)
btn_file = tk.Button(file_row, text="Tanlash", bg="#4a90e2", fg="white", font=("Arial", 10, "bold"), bd=0, relief="flat", command=choose_file)
btn_file.pack(side="left")
btn_file.bind("<Enter>", lambda e: btn_file.config(bg="#357ABD"))
btn_file.bind("<Leave>", lambda e: btn_file.config(bg="#4a90e2"))

# Yo'nalish
opts_row = tk.Frame(card, bg="#2c2c2c")
opts_row.pack(pady=12)
mode_var = tk.StringVar(value="latin_to_cyr")
tk.Radiobutton(opts_row, text="Lotin → Кирилл", variable=mode_var, value="latin_to_cyr", bg="#2c2c2c", fg="#ffffff", selectcolor="#4a90e2", font=("Arial", 11)).pack(side="left", padx=20)
tk.Radiobutton(opts_row, text="Кирилл → Lotin", variable=mode_var, value="cyr_to_latin", bg="#2c2c2c", fg="#ffffff", selectcolor="#4a90e2", font=("Arial", 11)).pack(side="left", padx=20)

# Action buttons
btn_row = tk.Frame(card, bg="#2c2c2c")
btn_row.pack(pady=15)
def style_btn(btn, color, hover_color):
    btn['bg']=color; btn['fg']="white"; btn['bd']=0; btn['relief']="flat"; btn['font']=("Arial", 11, "bold")
    btn.bind("<Enter>", lambda e: btn.config(bg=hover_color))
    btn.bind("<Leave>", lambda e: btn.config(bg=color))
def do_translate(mode):
    input_path = file_var.get()
    if not input_path:
        messagebox.showwarning("Xato", "Iltimos, PowerPoint faylni tanlang.")
        return
    base = os.path.splitext(input_path)[0]
    out = base + ("_krill.pptx" if mode=="latin_to_cyr" else "_lotin.pptx")
    try:
        process_pptx(input_path, out, direction=mode)
        messagebox.showinfo("Tayyor", f"Fayl saqlandi:\n{out}")
    except Exception as e:
        messagebox.showerror("Xato", f"Tarjima jarayonida xato:\n{e}")
btn_translate = tk.Button(btn_row, text="Tarjima qil", command=lambda: do_translate(mode_var.get()))
btn_translate.pack(side="left", padx=15, ipadx=10)
style_btn(btn_translate, "#27ae60", "#2ecc71")
btn_exit = tk.Button(btn_row, text="Chiqish", command=root.destroy)
btn_exit.pack(side="left", padx=15, ipadx=10)
style_btn(btn_exit, "#c0392b", "#e74c3c")

# Footer hints
hint_frame = tk.Frame(card, bg="#2c2c2c")
hint_frame.pack(fill="x", padx=20)
tk.Label(hint_frame, text="Eslatma: dastur .pptx formatini afzal ko‘radi. Agar .ppt bo‘lsa, PowerPointda ochib .pptx qilib saqlang.",
         font=("Arial", 9), fg="#cccccc", bg="#2c2c2c", wraplength=600, justify="left").pack(anchor="w", pady=(5,0))
tk.Label(hint_frame, text="Ayrim belgilarni va harflarni o‘zingiz mustaqil to‘g‘irlashingiz shart, chunki krillchada ayrim so'zlarni tarjima qilishda muammolar mavjud",
         font=("Arial", 9), fg="#cccccc", bg="#2c2c2c", wraplength=600, justify="left").pack(anchor="w", pady=(2,0))

root.mainloop()
