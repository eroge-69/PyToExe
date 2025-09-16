import sys
import os
from PyQt5.QtWidgets import (
    QApplication, QWidget, QTabWidget, QVBoxLayout, QHBoxLayout, QPushButton,
    QLabel, QTextEdit, QFileDialog, QProgressBar, QInputDialog, QMessageBox, QSizePolicy
)
from PyQt5.QtCore import Qt, QObject, QThread, pyqtSignal

# ---- Theme Stylesheets ----
LIGHT_THEME = """
QWidget { background: #fafafa; color: #222; }
QPushButton { background: #e0e0e0; color: #222; border-radius:6px; padding:5px 14px; }
QPushButton:hover { background: #d5d5ff; }
QTextEdit, QLineEdit { background: #fff; color: #111; }
QTabWidget::pane { border: 1px solid #bbb; }
QTabBar::tab { background: #e0e0e0; color: #222; padding: 6px; margin: 2px; }
QTabBar::tab:selected, QTabBar::tab:hover { background: #b0c4ff; color: #222; }
QProgressBar { background: #e0e0e0; border: 1px solid #aaa; text-align:center;}
QProgressBar::chunk{ background: #4285f4; }
"""

DARK_THEME = """
QWidget { background: #222; color: #fafafa; }
QPushButton { background: #444; color: #fafafa; border-radius:6px; padding:5px 14px; }
QPushButton:hover { background: #333366; }
QTextEdit, QLineEdit { background: #333; color: #eee; }
QTabWidget::pane { border: 1px solid #666; }
QTabBar::tab { background: #444; color: #fafafa; padding: 6px; margin: 2px; }
QTabBar::tab:selected, QTabBar::tab:hover { background: #448aff; color: #fff; }
QProgressBar { background: #444; border: 1px solid #888; text-align:center;}
QProgressBar::chunk{ background: #90caf9; }
"""

# ---- Worker Base ----
class BaseWorker(QObject):
    finished = pyqtSignal(str)
    error = pyqtSignal(str)
    progress = pyqtSignal(int)

from PyPDF2 import PdfReader, PdfWriter

# ---- Worker Classes ----
# (Same worker classes as before; all are copied here for completeness.)

class MergePDFsWorker(BaseWorker):
    def __init__(self, files, save_path):
        super().__init__()
        self.files = files
        self.save_path = save_path

    def run(self):
        try:
            writer = PdfWriter()
            total_pages = sum(len(PdfReader(f).pages) for f in self.files)
            done_pages = 0
            for f in self.files:
                reader = PdfReader(f)
                if reader.is_encrypted:
                    raise Exception(f"File {os.path.basename(f)} is password-protected. Unlock first.")
                for page in reader.pages:
                    writer.add_page(page)
                    done_pages += 1
                    self.progress.emit(int((done_pages / total_pages) * 100))
            with open(self.save_path, "wb") as f_out:
                writer.write(f_out)
            self.finished.emit(f"Merged PDFs saved to: {self.save_path}")
        except Exception as e:
            self.error.emit(f"Error (Merge): {str(e)}")

class SplitPDFWorker(BaseWorker):
    def __init__(self, pdf_path, save_dir):
        super().__init__()
        self.pdf_path = pdf_path
        self.save_dir = save_dir

    def run(self):
        try:
            reader = PdfReader(self.pdf_path)
            if reader.is_encrypted:
                raise Exception("PDF is password-protected. Unlock first.")
            total = len(reader.pages)
            for i, page in enumerate(reader.pages):
                writer = PdfWriter()
                writer.add_page(page)
                output_path = os.path.join(self.save_dir, f"page_{i+1}.pdf")
                with open(output_path, "wb") as out_f:
                    writer.write(out_f)
                self.progress.emit(int((i+1)/total*100))
            self.finished.emit("PDF split into individual pages!")
        except Exception as e:
            self.error.emit(f"Error (Split): {str(e)}")

class PDFToWordWorker(BaseWorker):
    def __init__(self, pdf_path, docx_path):
        super().__init__()
        self.pdf_path = pdf_path
        self.docx_path = docx_path

    def run(self):
        try:
            from pdf2docx import Converter
            cv = Converter(self.pdf_path)
            cv.convert(self.docx_path, start=0, end=None)
            cv.close()
            self.finished.emit(f"Converted PDF to Word: {self.docx_path}")
        except Exception as e:
            self.error.emit(f"Error (PDF to Word): {str(e)}")

class PDFToExcelWorker(BaseWorker):
    def __init__(self, pdf_path, excel_path):
        super().__init__()
        self.pdf_path = pdf_path
        self.excel_path = excel_path

    def run(self):
        try:
            import tabula
            tabula.convert_into(self.pdf_path, self.excel_path, output_format="xlsx")
            self.finished.emit(f"Converted PDF to Excel: {self.excel_path}")
        except ImportError:
            self.error.emit("tabula-py is required. Install with 'pip install tabula-py'. Java is also required.")
        except Exception as e:
            self.error.emit(f"Error (PDF to Excel): {str(e)}")

class PDFToPPTWorker(BaseWorker):
    def __init__(self, pdf_path, ppt_path):
        super().__init__()
        self.pdf_path = pdf_path
        self.ppt_path = ppt_path

    def run(self):
        try:
            from pdf2image import convert_from_path
            from pptx import Presentation
            import os
            slides = convert_from_path(self.pdf_path)
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]
            for i, slide_img in enumerate(slides):
                slide = prs.slides.add_slide(blank_slide_layout)
                img_path = f"temp_slide_{i}.png"
                slide_img.save(img_path, "PNG")
                slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
                os.remove(img_path)
                self.progress.emit(int((i+1)/len(slides)*100))
            prs.save(self.ppt_path)
            self.finished.emit(f"Converted PDF to PPT: {self.ppt_path}")
        except ImportError:
            self.error.emit("Requires 'pdf2image' and 'python-pptx'. Install with 'pip install pdf2image python-pptx'. Also install 'poppler' for pdf2image.")
        except Exception as e:
            self.error.emit(f"Error (PDF to PPT): {str(e)}")

class CompressPDFWorker(BaseWorker):
    def __init__(self, pdf_path, save_path):
        super().__init__()
        self.pdf_path = pdf_path
        self.save_path = save_path

    def run(self):
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(self.pdf_path)
            doc.save(self.save_path, deflate=True)
            self.finished.emit(f"Compressed PDF saved to: {self.save_path}")
        except ImportError:
            self.error.emit("PyMuPDF is required. Install with 'pip install pymupdf'.")
        except Exception as e:
            self.error.emit(f"Error (Compress PDF): {str(e)}")

class UnlockPDFWorker(BaseWorker):
    def __init__(self, pdf_path, password, save_path):
        super().__init__()
        self.pdf_path = pdf_path
        self.password = password
        self.save_path = save_path

    def run(self):
        try:
            reader = PdfReader(self.pdf_path)
            if reader.is_encrypted:
                if not reader.decrypt(self.password):
                    raise Exception("Wrong password")
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            with open(self.save_path, "wb") as f_out:
                writer.write(f_out)
            self.finished.emit(f"Unlocked PDF saved to: {self.save_path}")
        except Exception as e:
            self.error.emit(f"Error (Unlock PDF): {str(e)}")

class LockPDFWorker(BaseWorker):
    def __init__(self, pdf_path, password, save_path):
        super().__init__()
        self.pdf_path = pdf_path
        self.password = password
        self.save_path = save_path

    def run(self):
        try:
            reader = PdfReader(self.pdf_path)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            writer.encrypt(self.password)
            with open(self.save_path, "wb") as f_out:
                writer.write(f_out)
            self.finished.emit(f"Locked PDF saved to: {self.save_path}")
        except Exception as e:
            self.error.emit(f"Error (Lock PDF): {str(e)}")

class OrganizePDFWorker(BaseWorker):
    def __init__(self, pdf_path, order_list, save_path):
        super().__init__()
        self.pdf_path = pdf_path
        self.order_list = order_list
        self.save_path = save_path

    def run(self):
        try:
            reader = PdfReader(self.pdf_path)
            num_pages = len(reader.pages)
            writer = PdfWriter()
            for idx in self.order_list:
                if 0 <= idx < num_pages:
                    writer.add_page(reader.pages[idx])
                else:
                    raise Exception(f"Invalid page index: {idx+1}")
            with open(self.save_path, "wb") as f_out:
                writer.write(f_out)
            self.finished.emit(f"Organized PDF saved to: {self.save_path}")
        except Exception as e:
            self.error.emit(f"Error (Organize PDF): {str(e)}")

class WordToPDFWorker(BaseWorker):
    def __init__(self, docx_path, pdf_path):
        super().__init__()
        self.docx_path = docx_path
        self.pdf_path = pdf_path

    def run(self):
        try:
            from docx import Document
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            doc = Document(self.docx_path)
            c = canvas.Canvas(self.pdf_path, pagesize=letter)
            width, height = letter
            textobject = c.beginText(40, height - 40)
            for para in doc.paragraphs:
                textobject.textLine(para.text)
            c.drawText(textobject)
            c.save()
            self.finished.emit(f"Word to PDF conversion done: {self.pdf_path}")
        except Exception as e:
            self.error.emit(f"Error (Word to PDF): {str(e)}")

class ExcelToPDFWorker(BaseWorker):
    def __init__(self, xlsx_path, pdf_path):
        super().__init__()
        self.xlsx_path = xlsx_path
        self.pdf_path = pdf_path

    def run(self):
        try:
            import pandas as pd
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            df = pd.read_excel(self.xlsx_path)
            c = canvas.Canvas(self.pdf_path, pagesize=letter)
            width, height = letter
            textobject = c.beginText(40, height - 40)
            textobject.textLine("Excel to PDF Output:")
            for row in df.itertuples(index=False):
                textobject.textLine(', '.join(str(cell) for cell in row))
            c.drawText(textobject)
            c.save()
            self.finished.emit(f"Excel to PDF conversion done: {self.pdf_path}")
        except Exception as e:
            self.error.emit(f"Error (Excel to PDF): {str(e)}")

class PPTToPDFWorker(BaseWorker):
    def __init__(self, pptx_path, pdf_path):
        super().__init__()
        self.pptx_path = pptx_path
        self.pdf_path = pdf_path

    def run(self):
        try:
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            prs = Presentation(self.pptx_path)
            c = canvas.Canvas(self.pdf_path, pagesize=letter)
            width, height = letter
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
                textobject = c.beginText(40, height - 40)
                for line in texts:
                    textobject.textLine(line)
                c.drawText(textobject)
                c.showPage()
                self.progress.emit(int((i+1)/len(prs.slides)*100))
            c.save()
            self.finished.emit(f"PPT to PDF conversion done: {self.pdf_path}")
        except Exception as e:
            self.error.emit(f"Error (PPT to PDF): {str(e)}")

class OCRPDFWorker(BaseWorker):
    def __init__(self, pdf_path, txt_path):
        super().__init__()
        self.pdf_path = pdf_path
        self.txt_path = txt_path

    def run(self):
        try:
            from pdf2image import convert_from_path
            from pytesseract import image_to_string
            images = convert_from_path(self.pdf_path)
            output_text = ''
            for i, img in enumerate(images):
                text = image_to_string(img)
                output_text += f"--- Page {i+1} ---\n{text}\n"
                self.progress.emit(int((i+1)/len(images)*100))
            with open(self.txt_path, "w", encoding="utf-8") as f:
                f.write(output_text)
            self.finished.emit(f"OCR complete, text saved to: {self.txt_path}")
        except Exception as e:
            self.error.emit(f"Error (OCR PDF): {str(e)}")

class ScanPDFWorker(BaseWorker):
    def __init__(self, pdf_path, num_pages=1):
        super().__init__()
        self.pdf_path = pdf_path
        self.num_pages = num_pages

    def run(self):
        try:
            import cv2
            from PIL import Image
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            cap = cv2.VideoCapture(0)
            imgs = []
            for i in range(self.num_pages):
                ret, frame = cap.read()
                if not ret:
                    raise Exception("Failed to capture image from camera.")
                img = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                pil_img = Image.fromarray(img)
                img_path = f"scan_temp_{i}.png"
                pil_img.save(img_path)
                imgs.append(img_path)
                self.progress.emit(int((i+1)/self.num_pages*100))
            cap.release()
            c = canvas.Canvas(self.pdf_path, pagesize=letter)
            for img_path in imgs:
                c.drawImage(img_path, 0, 0, *letter)
                c.showPage()
                os.remove(img_path)
            c.save()
            self.finished.emit(f"Scan to PDF complete: {self.pdf_path}")
        except Exception as e:
            self.error.emit(f"Error (Scan PDF): {str(e)}")

class RedactPDFWorker(BaseWorker):
    def __init__(self, pdf_path, redact_text, save_path):
        super().__init__()
        self.pdf_path = pdf_path
        self.redact_text = redact_text
        self.save_path = save_path

    def run(self):
        try:
            reader = PdfReader(self.pdf_path)
            writer = PdfWriter()
            for i, page in enumerate(reader.pages):
                writer.add_page(page)
                self.progress.emit(int((i+1)/len(reader.pages)*100))
            with open(self.save_path, "wb") as f_out:
                writer.write(f_out)
            self.finished.emit(f"Redact complete (placeholder, not actual removal). Saved: {self.save_path}")
        except Exception as e:
            self.error.emit(f"Error (Redact PDF): {str(e)}")

# ---- Shared Header Bar for Tabs ----
class FeatureHeader(QWidget):
    theme_toggled = pyqtSignal()
    back_home = pyqtSignal()

    def __init__(self, feature_name):
        super().__init__()
        box = QHBoxLayout()
        fname = QLabel(f"Feature: {feature_name}")
        fname.setStyleSheet("font-size:18px; font-weight:bold;")
        box.addWidget(fname)
        box.addStretch()
        theme_btn = QPushButton("🌗 Toggle Theme")
        theme_btn.clicked.connect(self.theme_toggled.emit)
        box.addWidget(theme_btn)
        home_btn = QPushButton("🏠 Back to Home")
        home_btn.clicked.connect(self.back_home.emit)
        box.addWidget(home_btn)
        self.setLayout(box)

# ---- Feature Tab Template ----
class FeatureTab(QWidget):
    def __init__(self, feature_name, parent=None):
        super().__init__(parent)
        self.feature_name = feature_name
        self.header = FeatureHeader(feature_name)
        self.progress = QProgressBar()
        self.progress.setValue(0)
        self.status_box = QTextEdit()
        self.status_box.setReadOnly(True)
        self.vbox = QVBoxLayout()
        self.vbox.addWidget(self.header)
        self.vbox.addWidget(self.status_box)
        self.vbox.addWidget(self.progress)
        self.setLayout(self.vbox)

    def show_status(self, msg):
        self.status_box.append(msg)
    def reset_progress(self):
        self.progress.setValue(0)
    def set_progress(self, val):
        self.progress.setValue(val)

# ---- All Real Feature Tabs ----
# (For every tab, the pattern is the same: select file(s), run worker, connect progress/finished/error, show status!)

class MergeTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Merge PDF Files", parent)
        btn_select = QPushButton("Select PDF Files to Merge")
        btn_select.clicked.connect(self.select_files)
        self.vbox.insertWidget(1, btn_select)
        self.files = []
        self.save_path = ""
        self.worker = None
        self.thread = None

    def select_files(self):
        files, _ = QFileDialog.getOpenFileNames(self, "Select PDF Files", "", filter="PDF Files (*.pdf)")
        if not files: return
        self.files = files
        save_path, _ = QFileDialog.getSaveFileName(self, "Save Merged PDF", "", filter="PDF Files (*.pdf)")
        if not save_path: return
        self.save_path = save_path
        self.run_worker()

    def run_worker(self):
        self.show_status("Merging PDFs, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = MergePDFsWorker(self.files, self.save_path)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.progress.connect(self.progress.setValue)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

    def dropEvent(self, event):
        files = [u.toLocalFile() for u in event.mimeData().urls() if u.toLocalFile().lower().endswith('.pdf')]
        if files:
            self.files = files
            save_path, _ = QFileDialog.getSaveFileName(self, "Save Merged PDF", "", filter="PDF Files (*.pdf)")
            if not save_path: return
            self.save_path = save_path
            self.run_worker()

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()

class SplitTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Split PDF", parent)
        btn_select = QPushButton("Select PDF to Split")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.pdf = ""
        self.save_dir = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        pdf, _ = QFileDialog.getOpenFileName(self, "Select PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        save_dir = QFileDialog.getExistingDirectory(self, "Select Output Directory", "")
        if not save_dir: return
        self.save_dir = save_dir
        self.run_worker()

    def run_worker(self):
        self.show_status("Splitting PDF, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = SplitPDFWorker(self.pdf, self.save_dir)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.progress.connect(self.progress.setValue)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- PDF to Word ---
class PDFToWordTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Convert PDF to Word", parent)
        btn_select = QPushButton("Select PDF to Convert")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.pdf = ""
        self.docx = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        pdf, _ = QFileDialog.getOpenFileName(self, "Select PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        docx, _ = QFileDialog.getSaveFileName(self, "Save DOCX File", "", filter="Word Files (*.docx)")
        if not docx: return
        self.docx = docx
        self.run_worker()

    def run_worker(self):
        self.show_status("Converting PDF to Word, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = PDFToWordWorker(self.pdf, self.docx)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- PDF to Excel ---
class PDFToExcelTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Convert PDF to Excel", parent)
        btn_select = QPushButton("Select PDF to Convert")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.pdf = ""
        self.excel = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        pdf, _ = QFileDialog.getOpenFileName(self, "Select PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        excel, _ = QFileDialog.getSaveFileName(self, "Save Excel File", "", filter="Excel Files (*.xlsx)")
        if not excel: return
        self.excel = excel
        self.run_worker()

    def run_worker(self):
        self.show_status("Converting PDF to Excel, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = PDFToExcelWorker(self.pdf, self.excel)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- PDF to PPT ---
class PDFToPPTTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Convert PDF to PPT", parent)
        btn_select = QPushButton("Select PDF to Convert")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.pdf = ""
        self.ppt = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        pdf, _ = QFileDialog.getOpenFileName(self, "Select PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        ppt, _ = QFileDialog.getSaveFileName(self, "Save PPT File", "", filter="PowerPoint Files (*.pptx)")
        if not ppt: return
        self.ppt = ppt
        self.run_worker()

    def run_worker(self):
        self.show_status("Converting PDF to PPT, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = PDFToPPTWorker(self.pdf, self.ppt)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.progress.connect(self.progress.setValue)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- Compress PDF ---
class CompressTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Compress PDF", parent)
        btn_select = QPushButton("Select PDF to Compress")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.pdf = ""
        self.save = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        pdf, _ = QFileDialog.getOpenFileName(self, "Select PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        save, _ = QFileDialog.getSaveFileName(self, "Save Compressed PDF", "", filter="PDF Files (*.pdf)")
        if not save: return
        self.save = save
        self.run_worker()

    def run_worker(self):
        self.show_status("Compressing PDF, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = CompressPDFWorker(self.pdf, self.save)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- Unlock PDF ---
class UnlockTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Unlock PDF", parent)
        btn_select = QPushButton("Select Encrypted PDF to Unlock")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.pdf = ""
        self.save = ""
        self.password = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        pdf, _ = QFileDialog.getOpenFileName(self, "Select Encrypted PDF", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        password, ok = QInputDialog.getText(self, 'Unlock PDF', 'Enter password:')
        if not ok: return
        self.password = password
        save, _ = QFileDialog.getSaveFileName(self, "Save Unlocked PDF", "", filter="PDF Files (*.pdf)")
        if not save: return
        self.save = save
        self.run_worker()

    def run_worker(self):
        self.show_status("Unlocking PDF, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = UnlockPDFWorker(self.pdf, self.password, self.save)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- Lock PDF ---
class LockTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Lock PDF", parent)
        btn_select = QPushButton("Select PDF to Lock")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.pdf = ""
        self.save = ""
        self.password = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        pdf, _ = QFileDialog.getOpenFileName(self, "Select PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        password, ok = QInputDialog.getText(self, 'Lock PDF', 'Enter new password:')
        if not ok: return
        self.password = password
        save, _ = QFileDialog.getSaveFileName(self, "Save Locked PDF", "", filter="PDF Files (*.pdf)")
        if not save: return
        self.save = save
        self.run_worker()

    def run_worker(self):
        self.show_status("Locking PDF, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = LockPDFWorker(self.pdf, self.password, self.save)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- Organize PDF ---
class OrganizeTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Organize PDF", parent)
        btn_select = QPushButton("Select PDF to Organize")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.pdf = ""
        self.save = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        pdf, _ = QFileDialog.getOpenFileName(self, "Select PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        reader = PdfReader(pdf)
        num_pages = len(reader.pages)
        order, ok = QInputDialog.getText(self, 'Organize PDF', f'Enter new page order (e.g., 2,1,3) for {num_pages} pages:')
        if not ok: return
        try:
            order_list = [int(x.strip())-1 for x in order.split(",") if x.strip().isdigit()]
        except Exception:
            self.show_status("Invalid input for page order!")
            return
        save, _ = QFileDialog.getSaveFileName(self, "Save Organized PDF", "", filter="PDF Files (*.pdf)")
        if not save: return
        self.save = save
        self.run_worker(order_list)

    def run_worker(self, order_list):
        self.show_status("Organizing PDF, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = OrganizePDFWorker(self.pdf, order_list, self.save)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- Word to PDF ---
class WordToPDFTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Word to PDF", parent)
        btn_select = QPushButton("Select Word File to Convert")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.docx = ""
        self.pdf = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        docx, _ = QFileDialog.getOpenFileName(self, "Select Word File", "", filter="Word Files (*.docx)")
        if not docx: return
        self.docx = docx
        pdf, _ = QFileDialog.getSaveFileName(self, "Save PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        self.run_worker()

    def run_worker(self):
        self.show_status("Converting Word to PDF, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = WordToPDFWorker(self.docx, self.pdf)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- Excel to PDF ---
class ExcelToPDFTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Excel to PDF", parent)
        btn_select = QPushButton("Select Excel File to Convert")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.xlsx = ""
        self.pdf = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        xlsx, _ = QFileDialog.getOpenFileName(self, "Select Excel File", "", filter="Excel Files (*.xlsx)")
        if not xlsx: return
        self.xlsx = xlsx
        pdf, _ = QFileDialog.getSaveFileName(self, "Save PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        self.run_worker()

    def run_worker(self):
        self.show_status("Converting Excel to PDF, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = ExcelToPDFWorker(self.xlsx, self.pdf)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- PPT to PDF ---
class PPTToPDFTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("PPT to PDF", parent)
        btn_select = QPushButton("Select PPT File to Convert")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.pptx = ""
        self.pdf = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        pptx, _ = QFileDialog.getOpenFileName(self, "Select PPT File", "", filter="PowerPoint Files (*.pptx)")
        if not pptx: return
        self.pptx = pptx
        pdf, _ = QFileDialog.getSaveFileName(self, "Save PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        self.run_worker()

    def run_worker(self):
        self.show_status("Converting PPT to PDF, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = PPTToPDFWorker(self.pptx, self.pdf)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.progress.connect(self.progress.setValue)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- OCR PDF ---
class OCRTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("OCR PDF", parent)
        btn_select = QPushButton("Select PDF for OCR")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.pdf = ""
        self.txt = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        pdf, _ = QFileDialog.getOpenFileName(self, "Select PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        txt, _ = QFileDialog.getSaveFileName(self, "Save Text File", "", filter="Text Files (*.txt)")
        if not txt: return
        self.txt = txt
        self.run_worker()

    def run_worker(self):
        self.show_status("Running OCR on PDF, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = OCRPDFWorker(self.pdf, self.txt)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.progress.connect(self.progress.setValue)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- Scan PDF ---
class ScanTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Scan PDF", parent)
        btn_select = QPushButton("Scan pages to PDF from webcam")
        btn_select.clicked.connect(self.scan_pdf)
        self.vbox.insertWidget(1, btn_select)
        self.pdf = ""
        self.num_pages = 1
        self.worker = None
        self.thread = None

    def scan_pdf(self):
        pdf, _ = QFileDialog.getSaveFileName(self, "Save Scanned PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        num_pages, ok = QInputDialog.getInt(self, "Scan PDF", "Number of pages to scan:", 1, 1, 20, 1)
        if not ok: return
        self.num_pages = num_pages
        self.run_worker()

    def run_worker(self):
        self.show_status("Scanning pages, please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = ScanPDFWorker(self.pdf, self.num_pages)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.progress.connect(self.progress.setValue)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()

    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- Redact PDF ---
class RedactTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Redact PDF", parent)
        btn_select = QPushButton("Select PDF to Redact")
        btn_select.clicked.connect(self.select_file)
        self.vbox.insertWidget(1, btn_select)
        self.pdf = ""
        self.save = ""
        self.redact_text = ""
        self.worker = None
        self.thread = None

    def select_file(self):
        pdf, _ = QFileDialog.getOpenFileName(self, "Select PDF File", "", filter="PDF Files (*.pdf)")
        if not pdf: return
        self.pdf = pdf
        redact_text, ok = QInputDialog.getText(self, "Redact PDF", "Enter text to redact:")
        if not ok or not redact_text: return
        self.redact_text = redact_text
        save, _ = QFileDialog.getSaveFileName(self, "Save Redacted PDF", "", filter="PDF Files (*.pdf)")
        if not save: return
        self.save = save
        self.run_worker()

    def run_worker(self):
        self.show_status("Redacting PDF (placeholder: actual redaction not implemented), please wait...")
        self.progress.setValue(0)
        self.thread = QThread()
        self.worker = RedactPDFWorker(self.pdf, self.redact_text, self.save)
        self.worker.moveToThread(self.thread)
        self.thread.started.connect(self.worker.run)
        self.worker.progress.connect(self.progress.setValue)
        self.worker.finished.connect(self.worker_done)
        self.worker.error.connect(self.worker_done)
        self.worker.finished.connect(self.thread.quit)
        self.worker.error.connect(self.thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.worker.error.connect(self.worker.deleteLater)
        self.thread.finished.connect(self.thread.deleteLater)
        self.thread.start()
    def worker_done(self, msg):
        self.show_status(msg)
        self.progress.setValue(0)

# --- PDF to PDF/A (Stub) ---
class PDFATab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("PDF to PDF/A", parent)
        label = QLabel("PDF to PDF/A conversion is not implemented.<br>Use a commercial PDF/A tool.", self)
        label.setWordWrap(True)
        self.vbox.insertWidget(1, label)

# --- Edit PDF (Stub) ---
class EditTab(FeatureTab):
    def __init__(self, parent=None):
        super().__init__("Edit PDF", parent)
        label = QLabel("Advanced PDF editing is not implemented.<br>You can merge, split, and organize pages.", self)
        label.setWordWrap(True)
        self.vbox.insertWidget(1, label)

# ---- Home Page Widget ----
class HomePage(QWidget):
    feature_selected = pyqtSignal(str)

    def __init__(self, features, parent=None):
        super().__init__(parent)
        self.features = features
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout()
        label = QLabel("PDF Toolkit - Home")
        label.setAlignment(Qt.AlignHCenter)
        label.setStyleSheet("font-size: 28px; font-weight: bold; margin-top:18px; margin-bottom:22px;")
        layout.addWidget(label)
        grid = QHBoxLayout()
        col1 = QVBoxLayout()
        col2 = QVBoxLayout()
        for i, (name, _) in enumerate(self.features):
            btn = QPushButton(name)
            btn.setMinimumHeight(42)
            btn.setStyleSheet("font-size: 18px;")
            btn.clicked.connect(lambda _, n=name: self.feature_selected.emit(n))
            (col1 if i % 2 == 0 else col2).addWidget(btn)
        grid.addLayout(col1)
        grid.addSpacing(30)
        grid.addLayout(col2)
        layout.addLayout(grid)
        layout.addStretch()
        self.setLayout(layout)

# ---- Main Application ----
class PDFToolkit(QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("PDF Toolkit - All-in-One (Tabbed)")
        self.resize(1200, 750)
        self.tabs = QTabWidget()
        self.tabs.setTabsClosable(True)
        self.tabs.tabCloseRequested.connect(self.close_tab)
        self.feature_map = {
            "Merge PDF Files": MergeTab,
            "Split PDF": SplitTab,
            "Convert PDF to Word": PDFToWordTab,
            "Convert PDF to Excel": PDFToExcelTab,
            "Convert PDF to PPT": PDFToPPTTab,
            "Compress PDF": CompressTab,
            "Unlock PDF": UnlockTab,
            "Lock PDF": LockTab,
            "Organize PDF": OrganizeTab,
            "Word to PDF": WordToPDFTab,
            "Excel to PDF": ExcelToPDFTab,
            "PPT to PDF": PPTToPDFTab,
            "OCR PDF": OCRTab,
            "Scan PDF": ScanTab,
            "Redact PDF": RedactTab,
            "PDF to PDF/A": PDFATab,
            "Edit PDF": EditTab
        }
        self.features = list(self.feature_map.items())
        self.theme_is_dark = False
        self.init_ui()

    def init_ui(self):
        vbox = QVBoxLayout()
        self.tabs.clear()
        self.home = HomePage(self.features, self)
        self.home.feature_selected.connect(self.open_feature_tab)
        self.tabs.addTab(self.home, "🏠 Home")
        vbox.addWidget(self.tabs)
        self.setLayout(vbox)
        self.apply_theme()

    def apply_theme(self):
        self.setStyleSheet(DARK_THEME if self.theme_is_dark else LIGHT_THEME)

    def toggle_theme(self):
        self.theme_is_dark = not self.theme_is_dark
        self.apply_theme()

    def open_feature_tab(self, feature_name):
        for i in range(self.tabs.count()):
            if self.tabs.tabText(i) == feature_name:
                self.tabs.setCurrentIndex(i)
                return
        tab_class = self.feature_map[feature_name]
        tab = tab_class()
        tab.header.theme_toggled.connect(self.toggle_theme)
        tab.header.back_home.connect(self.go_home)
        idx = self.tabs.addTab(tab, feature_name)
        self.tabs.setCurrentIndex(idx)

    def go_home(self):
        self.tabs.setCurrentIndex(0)

    def close_tab(self, index):
        if index == 0: return
        self.tabs.removeTab(index)

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
    def dropEvent(self, event):
        files = [u.toLocalFile() for u in event.mimeData().urls()]
        if files and files[0].lower().endswith('.pdf'):
            self.open_feature_tab("Merge PDF Files")
            tab = self.tabs.currentWidget()
            if hasattr(tab, "dropEvent"):
                tab.files = files
                save_path, _ = QFileDialog.getSaveFileName(self, "Save Merged PDF", "", filter="PDF Files (*.pdf)")
                if not save_path: return
                tab.save_path = save_path
                tab.run_worker()

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = PDFToolkit()
    window.show()
    sys.exit(app.exec_())