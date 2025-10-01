#!/usr/bin/env python3
"""
Modern Dark Theme PowerPoint Text Remover with Scrollable Interface
A sleek, flat-design interface for removing text from PowerPoint files
"""

import customtkinter as ctk
from tkinter import filedialog, messagebox
import threading
import os
from pathlib import Path
from pptx import Presentation

# Configure appearance
ctk.set_appearance_mode("Dark")  # Dark theme
ctk.set_default_color_theme("blue")  # Blue color theme

class ScrollableFrame(ctk.CTkScrollableFrame):
    def __init__(self, master, **kwargs):
        super().__init__(master, **kwargs)
        
        # configure grid of scrollable frame
        self.grid_columnconfigure(0, weight=1)

class ModernPPTTextRemover(ctk.CTk):
    def __init__(self):
        super().__init__()
        
        self.title("PowerPoint Text Remover")
        self.geometry("800x600")
        self.minsize(700, 550)
        
        # Configure grid
        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(1, weight=1)
        
        self.setup_ui()
        
    def setup_ui(self):
        # Header
        self.header_frame = ctk.CTkFrame(self, corner_radius=0, height=120)
        self.header_frame.grid(row=0, column=0, sticky="ew", padx=0, pady=0)
        self.header_frame.grid_columnconfigure(0, weight=1)
        self.header_frame.grid_propagate(False)  # Prevent height changes
        
        # App title
        self.title_label = ctk.CTkLabel(
            self.header_frame,
            text="PowerPoint Text Remover",
            font=ctk.CTkFont(size=24, weight="bold"),
            text_color="#FF5100"
        )
        self.title_label.grid(row=0, column=0, pady=(20, 5), padx=20, sticky="w")
        
        # Subtitle
        self.subtitle_label = ctk.CTkLabel(
            self.header_frame,
            text="Remove all text from PowerPoint presentations while preserving formatting",
            font=ctk.CTkFont(size=12),
            text_color="gray70"
        )
        self.subtitle_label.grid(row=1, column=0, pady=(0, 20), padx=20, sticky="w")
        
        # Main scrollable content frame
        self.scrollable_frame = ScrollableFrame(
            self, 
            corner_radius=0,
            fg_color="transparent"
        )
        self.scrollable_frame.grid(row=1, column=0, sticky="nsew", padx=0, pady=0)
        self.scrollable_frame.grid_columnconfigure(0, weight=1)
        
        # Main content container inside scrollable frame
        self.main_container = ctk.CTkFrame(self.scrollable_frame, fg_color="transparent")
        self.main_container.grid(row=0, column=0, sticky="nsew", padx=20, pady=10)
        self.main_container.grid_columnconfigure(0, weight=1)
        
        # Single file section
        self.single_file_section = ctk.CTkFrame(self.main_container, corner_radius=15)
        self.single_file_section.grid(row=0, column=0, sticky="ew", padx=0, pady=10)
        self.single_file_section.grid_columnconfigure(1, weight=1)
        
        ctk.CTkLabel(
            self.single_file_section,
            text="📄 Process Single File",
            font=ctk.CTkFont(size=16, weight="bold")
        ).grid(row=0, column=0, columnspan=3, sticky="w", pady=(15, 10), padx=15)
        
        # File selection
        self.file_entry = ctk.CTkEntry(
            self.single_file_section,
            placeholder_text="Select a PowerPoint file...",
            height=40,
            font=ctk.CTkFont(size=12)
        )
        self.file_entry.grid(row=1, column=0, columnspan=2, sticky="ew", padx=(15, 10), pady=10)
        
        self.browse_file_btn = ctk.CTkButton(
            self.single_file_section,
            text="Browse",
            command=self.browse_file,
            width=100,
            height=40,
            fg_color="#FF5100",
            hover_color="#A44925",
            font=ctk.CTkFont(size=12, weight="bold")
        )
        self.browse_file_btn.grid(row=1, column=2, padx=(0, 15), pady=10)
        
        self.process_file_btn = ctk.CTkButton(
            self.single_file_section,
            text="Remove Text from File",
            command=self.process_single_file,
            height=45,
            fg_color="#FF00B3",
            hover_color="#B5368F",
            font=ctk.CTkFont(size=13, weight="bold")
        )
        self.process_file_btn.grid(row=2, column=0, columnspan=3, sticky="ew", padx=15, pady=(5, 15))
        
        # Batch processing section
        self.batch_section = ctk.CTkFrame(self.main_container, corner_radius=15)
        self.batch_section.grid(row=1, column=0, sticky="ew", padx=0, pady=10)
        self.batch_section.grid_columnconfigure(1, weight=1)
        
        ctk.CTkLabel(
            self.batch_section,
            text="📁 Process Multiple Files",
            font=ctk.CTkFont(size=16, weight="bold")
        ).grid(row=0, column=0, columnspan=3, sticky="w", pady=(15, 10), padx=15)
        
        # Directory selection
        self.dir_entry = ctk.CTkEntry(
            self.batch_section,
            placeholder_text="Select folder containing PowerPoint files...",
            height=40,
            font=ctk.CTkFont(size=12)
        )
        self.dir_entry.grid(row=1, column=0, columnspan=2, sticky="ew", padx=(15, 10), pady=10)
        
        self.browse_dir_btn = ctk.CTkButton(
            self.batch_section,
            text="Browse",
            command=self.browse_directory,
            width=100,
            height=40,
            fg_color="#FF5100",
            hover_color="#A44925",
            font=ctk.CTkFont(size=12, weight="bold")
        )
        self.browse_dir_btn.grid(row=1, column=2, padx=(0, 15), pady=10)
        
        self.process_batch_btn = ctk.CTkButton(
            self.batch_section,
            text="Remove Text from All Files",
            command=self.process_directory,
            height=45,
            fg_color="#FF00B3",
            hover_color="#B5368F",
            font=ctk.CTkFont(size=13, weight="bold")
        )
        self.process_batch_btn.grid(row=2, column=0, columnspan=3, sticky="ew", padx=15, pady=(5, 15))
        
        # Progress section
        self.progress_section = ctk.CTkFrame(self.main_container, corner_radius=15)
        self.progress_section.grid(row=2, column=0, sticky="ew", padx=0, pady=10)
        self.progress_section.grid_columnconfigure(0, weight=1)
        
        ctk.CTkLabel(
            self.progress_section,
            text="📊 Progress",
            font=ctk.CTkFont(size=16, weight="bold")
        ).grid(row=0, column=0, sticky="w", pady=(15, 10), padx=15)
        
        # Progress bar
        self.progress_bar = ctk.CTkProgressBar(
            self.progress_section, 
            height=20, 
            corner_radius=10,
            progress_color="#2CC985"
        )
        self.progress_bar.grid(row=1, column=0, sticky="ew", padx=15, pady=10)
        self.progress_bar.set(0)
        
        # Status label
        self.status_label = ctk.CTkLabel(
            self.progress_section,
            text="Ready to process files",
            font=ctk.CTkFont(size=12),
            text_color="gray70"
        )
        self.status_label.grid(row=2, column=0, sticky="w", pady=(0, 15), padx=15)
        
        # Results section
        self.results_frame = ctk.CTkFrame(self.main_container, corner_radius=15, fg_color="#1E3A5F")
        self.results_frame.grid(row=3, column=0, sticky="ew", padx=0, pady=10)
        self.results_frame.grid_columnconfigure(0, weight=1)
        
        self.results_label = ctk.CTkLabel(
            self.results_frame,
            text="📋 Processing Results",
            font=ctk.CTkFont(size=16, weight="bold"),
            text_color="#2CC985"
        )
        self.results_label.grid(row=0, column=0, sticky="w", pady=(15, 10), padx=15)
        
        self.results_text = ctk.CTkTextbox(
            self.results_frame,
            height=120,
            corner_radius=10,
            fg_color="#1A2B45",
            text_color="white",
            font=ctk.CTkFont(size=11),
            wrap="word"
        )
        self.results_text.grid(row=1, column=0, sticky="ew", padx=15, pady=(0, 15))
        self.results_text.insert("1.0", "Processing results will appear here...")
        self.results_text.configure(state="disabled")
        
        # Statistics section
        self.stats_frame = ctk.CTkFrame(self.main_container, corner_radius=15)
        self.stats_frame.grid(row=4, column=0, sticky="ew", padx=0, pady=10)
        self.stats_frame.grid_columnconfigure(0, weight=1)
        
        ctk.CTkLabel(
            self.stats_frame,
            text="📈 Statistics",
            font=ctk.CTkFont(size=16, weight="bold")
        ).grid(row=0, column=0, sticky="w", pady=(15, 15), padx=15)
        
        # Stats grid
        self.stats_grid = ctk.CTkFrame(self.stats_frame, fg_color="transparent")
        self.stats_grid.grid(row=1, column=0, sticky="ew", padx=15, pady=(0, 15))
        
        # Processed files stat
        self.processed_label = ctk.CTkLabel(
            self.stats_grid,
            text="Files Processed:",
            font=ctk.CTkFont(size=12),
            text_color="gray70"
        )
        self.processed_label.grid(row=0, column=0, sticky="w", padx=(0, 10))
        
        self.processed_value = ctk.CTkLabel(
            self.stats_grid,
            text="0",
            font=ctk.CTkFont(size=14, weight="bold"),
            text_color="#2CC985"
        )
        self.processed_value.grid(row=0, column=1, sticky="w", padx=(0, 30))
        
        # Text elements removed stat
        self.elements_label = ctk.CTkLabel(
            self.stats_grid,
            text="Text Elements Removed:",
            font=ctk.CTkFont(size=12),
            text_color="gray70"
        )
        self.elements_label.grid(row=0, column=2, sticky="w", padx=(0, 10))
        
        self.elements_value = ctk.CTkLabel(
            self.stats_grid,
            text="0",
            font=ctk.CTkFont(size=14, weight="bold"),
            text_color="#2CC985"
        )
        self.elements_value.grid(row=0, column=3, sticky="w")
        
        # Instructions section
        self.instructions_frame = ctk.CTkFrame(self.main_container, corner_radius=15)
        self.instructions_frame.grid(row=5, column=0, sticky="ew", padx=0, pady=10)
        
        ctk.CTkLabel(
            self.instructions_frame,
            text="💡 How to Use",
            font=ctk.CTkFont(size=16, weight="bold")
        ).grid(row=0, column=0, sticky="w", pady=(15, 10), padx=15)
        
        instructions_text = """• Single File: Select one PowerPoint file to remove all text
• Multiple Files: Select a folder to process all PowerPoint files
• Output: New files are created with '_no_text' suffix
• Original files are never modified
• Supports .pptx and .ppt file formats"""
        
        self.instructions_text = ctk.CTkTextbox(
            self.instructions_frame,
            height=100,
            corner_radius=10,
            fg_color="#1A2B45",
            text_color="gray80",
            font=ctk.CTkFont(size=11),
            wrap="word"
        )
        self.instructions_text.grid(row=1, column=0, sticky="ew", padx=15, pady=(0, 15))
        self.instructions_text.insert("1.0", instructions_text)
        self.instructions_text.configure(state="disabled")
        
        # Footer
        self.footer_frame = ctk.CTkFrame(self, corner_radius=0, fg_color="transparent")
        self.footer_frame.grid(row=2, column=0, sticky="ew", padx=20, pady=10)
        self.footer_frame.grid_columnconfigure(0, weight=1)
        
        ctk.CTkLabel(
            self.footer_frame,
            text="Made with ❤️ - PowerPoint Text Removal Tool v2.0",
            font=ctk.CTkFont(size=10),
            text_color="gray60"
        ).grid(row=0, column=0)
        
        # Initialize statistics
        self.processed_count = 0
        self.elements_count = 0
        
    def browse_file(self):
        filename = filedialog.askopenfilename(
            title="Select PowerPoint File",
            filetypes=[
                ("PowerPoint files", "*.pptx *.ppt"),
                ("PPTX files", "*.pptx"),
                ("PPT files", "*.ppt"),
                ("All files", "*.*")
            ]
        )
        if filename:
            self.file_entry.delete(0, "end")
            self.file_entry.insert(0, filename)
    
    def browse_directory(self):
        directory = filedialog.askdirectory(title="Select Folder with PowerPoint Files")
        if directory:
            self.dir_entry.delete(0, "end")
            self.dir_entry.insert(0, directory)
    
    def update_status(self, message, progress=None):
        self.status_label.configure(text=message)
        if progress is not None:
            self.progress_bar.set(progress)
        self.update_idletasks()
    
    def update_stats(self, processed_increment=0, elements_increment=0):
        self.processed_count += processed_increment
        self.elements_count += elements_increment
        self.processed_value.configure(text=str(self.processed_count))
        self.elements_value.configure(text=str(self.elements_count))
    
    def show_results(self, message, success=True):
        self.results_text.configure(state="normal")
        self.results_text.delete("1.0", "end")
        
        if success:
            self.results_label.configure(text="✅ Processing Complete!", text_color="#2CC985")
        else:
            self.results_label.configure(text="❌ Processing Failed", text_color="#FF6B6B")
            
        self.results_text.insert("1.0", message)
        self.results_text.configure(state="disabled")
        
        # Scroll to top of results
        self.results_text.see("1.0")
    
    def process_single_file(self):
        input_file = self.file_entry.get().strip()
        
        if not input_file:
            messagebox.showerror("Error", "Please select a PowerPoint file first.")
            return
        
        if not os.path.exists(input_file):
            messagebox.showerror("Error", "The selected file does not exist.")
            return
        
        # Disable buttons during processing
        self.process_file_btn.configure(state="disabled")
        self.process_batch_btn.configure(state="disabled")
        
        # Process in thread to keep UI responsive
        thread = threading.Thread(target=self._process_single_file_thread, args=(input_file,))
        thread.daemon = True
        thread.start()
    
    def _process_single_file_thread(self, input_file):
        try:
            self.update_status("Loading PowerPoint file...", 0.2)
            
            output_file = self._get_output_filename(input_file)
            result, elements_removed = remove_text_from_powerpoint(input_file, output_file)
            
            self.update_status("Complete!", 1.0)
            self.update_stats(processed_increment=1, elements_increment=elements_removed)
            
            # Show success message
            message = f"✅ Text removed successfully!\n\n" \
                     f"Input File: {os.path.basename(input_file)}\n" \
                     f"Output File: {os.path.basename(output_file)}\n" \
                     f"Text Elements Removed: {elements_removed}\n\n" \
                     f"Full Output Path:\n{output_file}"
            
            self.show_results(message, success=True)
            
            # Ask if user wants to open the output folder
            self.after(0, lambda: self.ask_open_folder(output_file))
            
        except Exception as e:
            error_msg = f"❌ Error processing file:\n\n{str(e)}\n\n" \
                       f"Please ensure:\n" \
                       f"• The file is not open in PowerPoint\n" \
                       f"• The file is a valid PowerPoint presentation\n" \
                       f"• You have permission to access the file"
            
            self.show_results(error_msg, success=False)
            self.update_status("Error occurred", 0)
        finally:
            # Re-enable buttons
            self.after(0, self.enable_buttons)
    
    def process_directory(self):
        directory = self.dir_entry.get().strip()
        
        if not directory:
            messagebox.showerror("Error", "Please select a directory first.")
            return
        
        if not os.path.exists(directory):
            messagebox.showerror("Error", "The selected directory does not exist.")
            return
        
        # Disable buttons during processing
        self.process_file_btn.configure(state="disabled")
        self.process_batch_btn.configure(state="disabled")
        
        # Process in thread
        thread = threading.Thread(target=self._process_directory_thread, args=(directory,))
        thread.daemon = True
        thread.start()
    
    def _process_directory_thread(self, directory):
        try:
            self.update_status("Scanning directory...", 0.1)
            
            input_path = Path(directory)
            ppt_files = list(input_path.glob("*.pptx")) + list(input_path.glob("*.ppt"))
            
            if not ppt_files:
                self.show_results("❌ No PowerPoint files found in the selected directory.\n\n"
                                 "Please ensure the directory contains .pptx or .ppt files.", success=False)
                self.update_status("No files found", 0)
                return
            
            self.update_status(f"Found {len(ppt_files)} files to process...", 0.3)
            
            output_dir = input_path / "text_free_output"
            output_dir.mkdir(exist_ok=True)
            
            processed_files = []
            failed_files = []
            total_elements_removed = 0
            
            for i, input_file in enumerate(ppt_files):
                progress = 0.3 + (i / len(ppt_files)) * 0.7
                self.update_status(f"Processing {i+1}/{len(ppt_files)}: {input_file.name}", progress)
                
                try:
                    output_file = output_dir / f"{input_file.stem}_no_text{input_file.suffix}"
                    result, elements_removed = remove_text_from_powerpoint(str(input_file), str(output_file))
                    processed_files.append(output_file)
                    total_elements_removed += elements_removed
                    
                    # Update stats after each file
                    self.update_stats(processed_increment=1, elements_increment=elements_removed)
                    
                except Exception as e:
                    failed_files.append(f"{input_file.name}: {str(e)}")
            
            # Prepare results message
            result_message = f"✅ Batch processing complete!\n\n"
            result_message += f"📊 Summary:\n"
            result_message += f"• Successfully processed: {len(processed_files)} files\n"
            result_message += f"• Failed: {len(failed_files)} files\n"
            result_message += f"• Total text elements removed: {total_elements_removed}\n\n"
            result_message += f"📁 Output folder: {output_dir}\n\n"
            
            if failed_files:
                result_message += f"❌ Failed files:\n" + "\n".join(failed_files)
            else:
                result_message += "🎉 All files processed successfully!"
            
            self.update_status("Complete!", 1.0)
            self.show_results(result_message, success=True)
            self.after(0, lambda: self.ask_open_folder(output_dir))
            
        except Exception as e:
            error_msg = f"❌ Error processing directory:\n\n{str(e)}\n\n" \
                       f"Please ensure:\n" \
                       f"• The directory is accessible\n" \
                       f"• You have permission to read/write files\n" \
                       f"• PowerPoint files are not currently open"
            
            self.show_results(error_msg, success=False)
            self.update_status("Error occurred", 0)
        finally:
            self.after(0, self.enable_buttons)
    
    def enable_buttons(self):
        self.process_file_btn.configure(state="normal")
        self.process_batch_btn.configure(state="normal")
    
    def _get_output_filename(self, input_file):
        input_path = Path(input_file)
        return str(input_path.parent / f"{input_path.stem}_no_text{input_path.suffix}")
    
    def ask_open_folder(self, file_or_folder_path):
        path = Path(file_or_folder_path)
        folder_path = path.parent if path.is_file() else path
        
        result = messagebox.askyesno(
            "Success", 
            "Processing complete! Would you like to open the output folder?"
        )
        
        if result:
            try:
                os.startfile(folder_path)  # Windows
            except:
                try:
                    import subprocess
                    subprocess.run(['open', folder_path])  # macOS
                except:
                    try:
                        subprocess.run(['xdg-open', folder_path])  # Linux
                    except:
                        messagebox.showinfo("Info", f"Output folder: {folder_path}")

def remove_text_from_powerpoint(input_file, output_file):
    """
    Remove all text from a PowerPoint presentation
    Returns: (output_file_path, elements_removed_count)
    """
    print(f"Processing: {input_file}")
    prs = Presentation(input_file)
    
    text_elements_removed = 0
    
    for slide_number, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            # Remove text from regular shapes
            if hasattr(shape, "text") and shape.text:
                original_text = shape.text
                shape.text = ""
                text_elements_removed += 1
                print(f"  Removed text from shape: '{original_text}'")
            
            # Handle tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            cell.text = ""
                            text_elements_removed += 1
            
            # Handle grouped shapes
            if shape.shape_type == 6:  # Group shape
                for sub_shape in shape.shapes:
                    if hasattr(sub_shape, "text") and sub_shape.text:
                        sub_shape.text = ""
                        text_elements_removed += 1
    
    prs.save(output_file)
    print(f"Saved text-free presentation: {output_file}")
    print(f"Total elements removed: {text_elements_removed}")
    return output_file, text_elements_removed

if __name__ == "__main__":
    app = ModernPPTTextRemover()
    app.mainloop()