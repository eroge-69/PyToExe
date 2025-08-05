
import os
import re
from docx import Document
import PyPDF2
from pptx import Presentation

def count_words_in_docx(file_path):
    doc = Document(file_path)
    text = ' '.join([p.text for p in doc.paragraphs])
    return len(re.findall(r'\w+', text))

def count_words_in_pdf(file_path):
    text = ''
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ''
    return len(re.findall(r'\w+', text))

def count_words_in_pptx(file_path):
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + ' '
    return len(re.findall(r'\w+', text))

def main(folder_path):
    total_words = 0
    results = []

    for root, dirs, files in os.walk(folder_path):
        for file in files:
            file_path = os.path.join(root, file)
            word_count = 0

            if file.lower().endswith('.docx'):
                word_count = count_words_in_docx(file_path)
            elif file.lower().endswith('.pdf'):
                word_count = count_words_in_pdf(file_path)
            elif file.lower().endswith('.pptx'):
                word_count = count_words_in_pptx(file_path)

            if word_count:
                results.append((file, word_count))
                total_words += word_count

    print('\nНәтиже (әр файл бойынша):')
    for file, count in results:
        print(f'{file}: {count} сөз')

    print(f'\nЖалпы сөз саны: {total_words} сөз')

if __name__ == "__main__":
    folder = input("Папканың толық жолын енгізіңіз (мысалы, C:\\Users\\User\\Documents\\myfolder): ")
    main(folder)
