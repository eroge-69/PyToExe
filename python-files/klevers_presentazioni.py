import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import zipfile, os

# === CONFIG ===
excel_file = "1.Potenziali Clienti Spagna (unified).xlsx"
logo_file = "klevers-logo-green-square.pdf"  # se disponibile in PNG meglio
output_zip = "Klevers_Presentazioni.zip"

# Palette colori dal logo Klevers (approssimata)
klevers_green = RGBColor(0, 122, 79)
klevers_dark = RGBColor(0, 51, 34)

# === Carica dati ===
df = pd.read_excel(excel_file, header=1)
df.columns = ["Index", "Settore", "Ragione Sociale", "Regione", "Provincia", "Email"]
df = df.drop(0).reset_index(drop=True)

# Statistiche
n_aziende = len(df)
settori = df["Settore"].value_counts()
regioni = df["Regione"].value_counts()

# === Funzione per creare slide titolo ===
def add_title_slide(prs, titolo):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = titolo
    subtitle.text = "Klevers Italiana Srl"
    return slide

# === Funzione per slide con grafico ===
def add_chart_slide(prs, titolo, data, chart_type="bar"):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title = shapes.title if slide.shapes.title else slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1)).text_frame
    title.text = titolo
    
    # Salva grafico come immagine temporanea
    fig, ax = plt.subplots()
    if chart_type == "bar":
        data.plot(kind="bar", ax=ax)
    else:
        data.plot(kind="pie", ax=ax, autopct='%1.1f%%')
    plt.title("")
    img_path = "temp_chart.png"
    plt.savefig(img_path)
    plt.close()
    
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), Inches(6), Inches(4))
    return slide

# === Funzione per creare presentazione scouting ===
def crea_presentazione_scouting(nome_file, stile="minimal"):
    prs = Presentation()
    add_title_slide(prs, "Scouting potenziali clienti Spagna")
    
    # Metodologia
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Metodologia di ricerca"
    tf = slide.placeholders[1].text_frame
    tf.text = "Google, Microsoft Edge, Europages"

    # Numero aziende
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Totale aziende censite"
    slide.placeholders[1].text = f"{n_aziende} aziende individuate"

    # Distribuzione per settore
    add_chart_slide(prs, "Distribuzione per settore", settori, chart_type="bar")
    add_chart_slide(prs, "Distribuzione per settore (torta)", settori, chart_type="pie")

    # Distribuzione per regione
    add_chart_slide(prs, "Distribuzione per regione", regioni, chart_type="bar")
    add_chart_slide(prs, "Distribuzione per regione (torta)", regioni, chart_type="pie")

    # Esempi aziende
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Esempi di aziende individuate"
    tf = slide.placeholders[1].text_frame
    for i, row in df.head(5).iterrows():
        tf.add_paragraph().text = f"{row['Ragione Sociale']} – {row['Email']}"

    # Conclusioni
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusioni strategiche"
    slide.placeholders[1].text = "Settori promettenti, opportunità B2B, contatti prioritari."

    # Raccomandazioni operative
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Raccomandazioni operative"
    tf = slide.placeholders[1].text_frame
    tf.text = ""
    for r in [
        "Approfondire i settori con più aziende censite",
        "Contattare i top prospect via email/LinkedIn",
        "Valutare partnership locali in Spagna",
        "Creare un database CRM interno",
        "Organizzare call di presentazione e follow-up"
    ]:
        tf.add_paragraph().text = r

    prs.save(nome_file)

# === Funzione per creare template ===
def crea_template(nome_file, stile="minimal"):
    prs = Presentation()
    add_title_slide(prs, "Template Klevers")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Esempio layout"
    slide.placeholders[1].text = "Testo segnaposto"
    prs.save(nome_file)

# === Crea i 4 pptx ===
files = []
crea_template("Template_Klevers_Minimal.pptx", "minimal")
files.append("Template_Klevers_Minimal.pptx")
crea_template("Template_Klevers_Dinamico.pptx", "dinamico")
files.append("Template_Klevers_Dinamico.pptx")
crea_presentazione_scouting("Scouting_Spagna_Minimal.pptx", "minimal")
files.append("Scouting_Spagna_Minimal.pptx")
crea_presentazione_scouting("Scouting_Spagna_Dinamico.pptx", "dinamico")
files.append("Scouting_Spagna_Dinamico.pptx")

# === Impacchetta in zip ===
with zipfile.ZipFile(output_zip, 'w') as zf:
    for f in files:
        zf.write(f)

print(f"Creato archivio {output_zip} con i 4 PowerPoint.")
